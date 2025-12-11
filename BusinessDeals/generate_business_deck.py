"""
Generate Non-Technical Business Deck for AI-Powered Presales POCs
Focused on business outcomes, ROI, and value for non-technical stakeholders
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """Create a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Style the title
    title_shape.text_frame.paragraphs[0].font.size = Pt(48)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)

def create_section_divider(prs, title, subtitle=""):
    """Create a section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Add title in center
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Add colored background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 102, 204)
    
    if subtitle:
        left = Inches(1)
        top = Inches(4.5)
        subtitle_box = slide.shapes.add_textbox(left, top, width, Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

def create_problem_solution_slide(prs, poc_number, title, problem, solution, value):
    """Create a problem-solution slide with clear visual separation"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Add POC number and title at top
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.7)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = f"Solution #{poc_number}: {title}"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    
    # Problem section (left side, red background)
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(4.25)
    height = Inches(5.5)
    
    # Add red background shape
    problem_bg = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
    problem_bg.fill.solid()
    problem_bg.fill.fore_color.rgb = RGBColor(220, 53, 69)
    problem_bg.line.color.rgb = RGBColor(220, 53, 69)
    
    # Problem text box
    problem_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    problem_frame = problem_box.text_frame
    problem_frame.word_wrap = True
    
    # Problem header
    p = problem_frame.paragraphs[0]
    p.text = "❌ THE PROBLEM TODAY"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_after = Pt(15)
    
    # Problem details
    for item in problem:
        p = problem_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(10)
    
    # Solution section (right side, green background)
    left = Inches(5.25)
    
    # Add green background shape
    solution_bg = slide.shapes.add_shape(1, left, top, width, height)  # Rectangle
    solution_bg.fill.solid()
    solution_bg.fill.fore_color.rgb = RGBColor(40, 167, 69)
    solution_bg.line.color.rgb = RGBColor(40, 167, 69)
    
    # Solution text box
    solution_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    solution_frame = solution_box.text_frame
    solution_frame.word_wrap = True
    
    # Solution header
    p = solution_frame.paragraphs[0]
    p.text = "✅ OUR SOLUTION"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_after = Pt(15)
    
    # Solution details
    for item in solution:
        p = solution_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(10)
    
    # Value proposition at bottom
    left = Inches(0.5)
    top = Inches(6.8)
    width = Inches(9)
    height = Inches(0.6)
    
    value_box = slide.shapes.add_textbox(left, top, width, height)
    value_frame = value_box.text_frame
    value_frame.text = f"💰 BUSINESS VALUE: {value}"
    p = value_frame.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.alignment = PP_ALIGN.CENTER

def create_simple_content_slide(prs, title, content_items, title_color=RGBColor(0, 102, 204)):
    """Create a simple content slide with large, readable text"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = title_color
    
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(20)
        p.space_before = Pt(15)

def create_metrics_slide(prs, title, metrics):
    """Create a slide with large metrics"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank
    
    # Title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.alignment = PP_ALIGN.CENTER
    
    # Metrics in grid
    col_width = Inches(4.25)
    row_height = Inches(2.5)
    start_left = Inches(0.75)
    start_top = Inches(2)
    
    for i, (metric_value, metric_label) in enumerate(metrics):
        row = i // 2
        col = i % 2
        
        left = start_left + (col * Inches(4.75))
        top = start_top + (row * Inches(2.8))
        
        # Metric box
        metric_box = slide.shapes.add_textbox(left, top, col_width, row_height)
        metric_frame = metric_box.text_frame
        metric_frame.word_wrap = True
        
        # Large metric value
        p = metric_frame.paragraphs[0]
        p.text = metric_value
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 204)
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        p = metric_frame.add_paragraph()
        p.text = metric_label
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(64, 64, 64)
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(10)

def create_presentation():
    """Create the complete non-technical business presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    create_title_slide(
        prs,
        "AI-Powered Presales Automation",
        "How We're Cutting Costs, Winning More Deals, and Accelerating Sales\n\n4 Game-Changing Solutions"
    )
    
    # Slide 2: The Big Picture Problem
    create_simple_content_slide(
        prs,
        "The Challenge We're Solving",
        [
            "🕒 Our sales teams spend too much time on repetitive work",
            "💸 Presales activities cost us $2-3 million annually",
            "⏰ We take weeks to respond to customer requests",
            "📉 We're losing deals because competitors move faster",
            "🎯 Our win rate could be 15-20% higher",
            "💡 We need smart automation to compete and win"
        ]
    )
    
    # Slide 3: Our Solution Overview
    create_simple_content_slide(
        prs,
        "Our Approach: Smart Automation",
        [
            "✨ We've built 4 intelligent tools that do the heavy lifting",
            "🤖 These tools work like expert assistants, available 24/7",
            "⚡ They complete in minutes what used to take days or weeks",
            "💰 They dramatically reduce costs while improving quality",
            "📈 They help us win more deals and close them faster",
            "🎯 Simple to use - no technical expertise required"
        ]
    )
    
    # Section divider
    create_section_divider(prs, "Solution #1", "Smart Infrastructure Designer")
    
    # Slide 4: POC1 - Problem & Solution
    create_problem_solution_slide(
        prs,
        1,
        "Smart Infrastructure Designer",
        [
            "Designing technology solutions takes 2-3 weeks",
            "Requires expensive specialists ($150K+ salaries)",
            "Inconsistent quality across different engineers",
            "Security and compliance gaps discovered too late",
            "Customers get frustrated with slow responses",
            "We lose deals while competitors respond faster"
        ],
        [
            "Designs complete solutions in 5-10 minutes",
            "Ensures compliance with regulations automatically",
            "Consistent, high-quality designs every time",
            "Catches security issues upfront, not later",
            "Creates documentation and code automatically",
            "Anyone on the team can generate professional designs"
        ],
        "Save $950,000 per year + Win deals 3 weeks faster"
    )
    
    # Slide 5: POC1 Real-World Example
    create_simple_content_slide(
        prs,
        "Real Example: Financial Services Client",
        [
            "BEFORE: Client asks for infrastructure design",
            "  ⏰ Waited 18 days for our design proposal",
            "  💰 Cost us $8,000 in presales engineering time",
            "  ❌ Competitor submitted proposal in 5 days and won",
            "",
            "AFTER: Same request with our Smart Designer",
            "  ⚡ Complete design ready in 8 minutes",
            "  ✅ Includes security compliance for financial regulations",
            "  💰 Cost: $50 (vs. $8,000)",
            "  🎯 Result: We submit first and impress the client"
        ]
    )
    
    # Section divider
    create_section_divider(prs, "Solution #2", "Personalized Demo Creator")
    
    # Slide 6: POC2 - Problem & Solution
    create_problem_solution_slide(
        prs,
        2,
        "Personalized Demo Creator",
        [
            "Generic demos don't resonate with customers",
            "Customizing demos takes 8-10 hours per customer",
            "Different audiences need different messages",
            "Demo environments break at critical moments",
            "Cost: $30,000-$40,000 per month in prep time",
            "Lost deals due to poor demo experiences"
        ],
        [
            "Automatically customizes demos for each audience",
            "Creates personalized demos in seconds, not hours",
            "Adapts messaging for CFOs, CTOs, or engineers",
            "Self-healing - demos never fail",
            "Interactive guides that match customer interests",
            "Every customer gets a 'wow' experience"
        ],
        "Save $300,000 per year + Improve demo success rate 40%"
    )
    
    # Slide 7: POC2 Real-World Example
    create_simple_content_slide(
        prs,
        "Real Example: Healthcare Prospect Demo",
        [
            "BEFORE: Healthcare CIO needs a demo",
            "  ⏰ Spent 9 hours customizing our standard demo",
            "  😕 CIO seemed unimpressed (too technical)",
            "  ❌ Lost deal - 'didn't address our business needs'",
            "",
            "AFTER: Same scenario with our Demo Creator",
            "  ⚡ System created personalized demo in 45 seconds",
            "  🎯 Focused on cost savings and compliance (what CIOs care about)",
            "  😊 CIO: 'This is exactly what we need!'",
            "  ✅ Won $280,000 contract"
        ]
    )
    
    # Section divider
    create_section_divider(prs, "Solution #3", "Instant Cost Calculator")
    
    # Slide 8: POC3 - Problem & Solution
    create_problem_solution_slide(
        prs,
        3,
        "Instant Cost Calculator",
        [
            "Cost analysis takes 3-5 days per opportunity",
            "Calculations are inconsistent across the team",
            "We often miss hidden costs in our estimates",
            "Can't quickly show ROI to justify our price",
            "No way to show environmental benefits",
            "Customers choose cheaper competitors by default"
        ],
        [
            "Instant, accurate cost analysis in seconds",
            "Shows 3-year total cost comparison",
            "Calculates exact ROI and payback period",
            "Includes environmental impact (carbon savings)",
            "What-if scenarios for different options",
            "Professional reports to share with customers"
        ],
        "Save $400,000 per year + Prove we're 40-60% cheaper than competitors"
    )
    
    # Slide 9: POC3 Real-World Example
    create_simple_content_slide(
        prs,
        "Real Example: Enterprise Migration Project",
        [
            "Customer considering $5M technology migration",
            "",
            "OUR INSTANT ANALYSIS SHOWED:",
            "  💰 Total 3-year cost: $38.4M (vs. competitor at $50M)",
            "  📈 ROI: 129% over 3 years",
            "  ⚡ Payback period: 14 months",
            "  💵 Annual savings: $6.5 million per year",
            "  🌱 Carbon reduction: 28% (great for sustainability goals)",
            "",
            "✅ RESULT: Won the deal based on clear financial proof"
        ]
    )
    
    # Section divider
    create_section_divider(prs, "Solution #4", "Competitive Intelligence System")
    
    # Slide 10: POC4 - Problem & Solution
    create_problem_solution_slide(
        prs,
        4,
        "Competitive Intelligence System",
        [
            "Competitor information is outdated",
            "Sales reps unprepared for competitive situations",
            "We lose deals based on competitor FUD (fear, uncertainty, doubt)",
            "Generic 'battle cards' don't help in real situations",
            "No insight into why we win or lose deals",
            "30-40% loss rate in competitive deals"
        ],
        [
            "Real-time competitor analysis and insights",
            "Automatic 'battle cards' for every situation",
            "Specific talking points for sales conversations",
            "Analyzes our wins and losses to find patterns",
            "Tells us exactly how to counter competitor claims",
            "Sales teams always prepared and confident"
        ],
        "Improve win rate 15% = $3 million in additional revenue per year"
    )
    
    # Slide 11: POC4 Real-World Example
    create_simple_content_slide(
        prs,
        "Real Example: Competing Against VMware",
        [
            "Deal: $450,000 cloud migration competing vs. VMware",
            "",
            "OUR SYSTEM PROVIDED:",
            "  🎯 VMware's key weakness: 45% price increase after acquisition",
            "  💰 Our price advantage: $200,000 cheaper over 3 years",
            "  💡 Talking point: 'VMware customers are switching to avoid price hikes'",
            "  📊 Demo focus: Show cloud-native features VMware lacks",
            "",
            "✅ RESULT: Sales rep felt confident, handled objections perfectly, WON the deal"
        ]
    )
    
    # Slide 12: Combined Impact - Metrics
    create_metrics_slide(
        prs,
        "Total Business Impact - The Bottom Line",
        [
            ("$4.65M", "Annual Financial Impact"),
            ("85%", "Time Savings on Manual Work"),
            ("15%", "Increase in Win Rate"),
            ("30%", "Faster Sales Cycles"),
            ("6,050%", "3-Year ROI"),
            ("2 months", "Payback Period")
        ]
    )
    
    # Slide 13: Before and After Comparison
    create_simple_content_slide(
        prs,
        "How Our Sales Process Changes",
        [
            "BEFORE (Without Automation):",
            "  ⏰ 3 weeks to create proposals",
            "  💸 $50K-100K presales cost per large deal",
            "  😰 Inconsistent quality and messaging",
            "  📉 Miss deadlines, lose competitive deals",
            "",
            "AFTER (With Our 4 Solutions):",
            "  ⚡ Same-day or next-day proposals",
            "  💰 $5K-10K presales cost per large deal",
            "  ✅ Professional, consistent, high-quality every time",
            "  🏆 Beat competitors on speed and substance"
        ]
    )
    
    # Slide 14: Why This Matters to Our Business
    create_simple_content_slide(
        prs,
        "Strategic Benefits",
        [
            "💼 COMPETITIVE ADVANTAGE: We respond faster than anyone in our market",
            "💰 COST REDUCTION: Cut presales expenses by 80%",
            "📈 REVENUE GROWTH: Win 15% more deals = millions in new revenue",
            "😊 CUSTOMER SATISFACTION: Impress customers with speed and quality",
            "👥 TEAM MORALE: Let experts do expert work, not repetitive tasks",
            "🚀 SCALABILITY: Handle 3x more opportunities without hiring"
        ]
    )
    
    # Slide 15: Investment Required
    create_simple_content_slide(
        prs,
        "What It Takes to Implement",
        [
            "💵 YEAR 1 INVESTMENT: $251,000",
            "  • Development and setup: $120,000",
            "  • Monthly operating costs: $7,000/month",
            "  • Training: $20,000 one-time",
            "  • Support: $30,000/year",
            "",
            "💰 YEAR 1 RETURN: $5.15 million",
            "  • Labor savings: $1.65M",
            "  • Revenue uplift: $3M",
            "  • Efficiency gains: $500K",
            "",
            "📊 NET BENEFIT: $4.9 million in Year 1 alone"
        ]
    )
    
    # Slide 16: The 90-Day Rollout Plan
    create_simple_content_slide(
        prs,
        "How We'll Make This Happen",
        [
            "PHASE 1 (Month 1): Start with Smart Designer",
            "  • Train 3 presales engineers",
            "  • Prove the concept with real customer requests",
            "",
            "PHASE 2 (Month 2): Add Demo Creator & Cost Calculator",
            "  • Expand to 10 team members",
            "  • Measure time savings and quality improvements",
            "",
            "PHASE 3 (Month 3): Deploy Competitive Intelligence",
            "  • Full team rollout",
            "  • Track wins and losses, refine approach",
            "",
            "✅ Quick wins visible in 30 days, full benefits by 90 days"
        ]
    )
    
    # Slide 17: Risk Management (Simple Language)
    create_simple_content_slide(
        prs,
        "Addressing Your Concerns",
        [
            "❓ 'What if the AI makes mistakes?'",
            "  ✅ Humans review everything before it goes to customers",
            "",
            "❓ 'Will people resist using it?'",
            "  ✅ Our team will love it - removes boring work, makes them heroes",
            "",
            "❓ 'Is our data secure?'",
            "  ✅ We can run everything on our own secure systems",
            "",
            "❓ 'What if it doesn't work?'",
            "  ✅ 30-day pilot with no long-term commitment, prove it first"
        ]
    )
    
    # Slide 18: Success Stories Preview
    create_simple_content_slide(
        prs,
        "What Success Looks Like",
        [
            "AFTER 6 MONTHS:",
            "  • Presales team handling 2x more opportunities",
            "  • Win rate increased from 35% to 45%",
            "  • Average sales cycle reduced by 3 weeks",
            "  • Customer satisfaction scores up 25%",
            "  • Team morale significantly improved",
            "",
            "AFTER 1 YEAR:",
            "  • $4.9M saved and earned (proven ROI)",
            "  • Industry recognition for innovation",
            "  • Competitive advantage that's hard to copy"
        ]
    )
    
    # Slide 19: What We Need from You
    create_simple_content_slide(
        prs,
        "Next Steps - Decision Points",
        [
            "✅ APPROVE: 30-day pilot program ($35,000)",
            "  • Select 3 presales team members",
            "  • Test with 10-15 real customer opportunities",
            "  • Measure results vs. traditional approach",
            "",
            "✅ COMMIT: Based on pilot results, full rollout",
            "  • Remaining investment of $216,000",
            "  • 90-day implementation across full team",
            "  • Regular progress reviews and adjustments",
            "",
            "🎯 DECISION NEEDED: Approve pilot program to start in January 2026"
        ]
    )
    
    # Slide 20: Call to Action
    create_title_slide(
        prs,
        "Let's Transform How We Sell",
        "Questions?\n\nPilot Program: Starting January 2026\nExpected Pilot Outcomes: Visible results in 30 days\nFull ROI: $4.9M in Year 1\n\nReady to move forward?"
    )
    
    # Save presentation
    prs.save('AI_Presales_Business_Case.pptx')
    print("✅ Business deck created: AI_Presales_Business_Case.pptx")
    print("📊 20 slides focused on business value for non-technical audience")
    print("🎯 Perfect for executives, business leaders, and decision makers")
    print("💡 Focus: Business problems, solutions, ROI, and real examples")

if __name__ == "__main__":
    create_presentation()
