"""
Azure Cost Management Orchestrator
Coordinates cost analysis and optimization, generates PowerPoint reports
"""

import os
import json
from datetime import datetime
from typing import Dict, List
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

# Import our analysis modules
from azure_cost_assessment import AzureCostAssessment
from cost_optimizer import AzureCostOptimizer
from assessment_config import get_config

# Load environment variables
load_dotenv()


class CostManagementOrchestrator:
    """Orchestrates cost analysis and optimization reporting"""
    
    def __init__(self):
        """Initialize the orchestrator"""
        self.config = get_config()
        self.cost_analysis_results = None
        self.optimization_results = None
        self.timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        
        print("🎯 Azure Cost Management Orchestrator")
        print("=" * 100)
    
    def run_cost_analysis(self, days: int = 30):
        """Run cost analysis"""
        print(f"\n📊 STEP 1: Running Cost Analysis ({days} days)...")
        print("-" * 100)
        
        try:
            assessment = AzureCostAssessment(self.config)
            
            # Get cost data
            cost_data = assessment.get_cost_data(days)
            
            if not cost_data:
                print("⚠️  No cost data available for analysis")
                self.cost_analysis_results = {
                    'total_cost': 0,
                    'daily_average': 0,
                    'by_service': [],
                    'by_resource_group': []
                }
                return
            
            # Analyze costs
            analysis = assessment.analyze_costs(cost_data)
            
            # Generate recommendations
            recommendations = assessment.generate_recommendations(analysis)
            
            self.cost_analysis_results = {
                'analysis': analysis,
                'recommendations': recommendations,
                'days_analyzed': days
            }
            
            print(f"✅ Cost Analysis Complete")
            print(f"   Total Cost: ${analysis['total_cost']:.2f}")
            print(f"   Daily Average: ${analysis['daily_average']:.2f}")
            print(f"   Top Services: {len(analysis['by_service'])}")
            
        except Exception as e:
            print(f"❌ Error in cost analysis: {e}")
            self.cost_analysis_results = None
    
    def run_optimization_analysis(self):
        """Run real-time optimization analysis"""
        print(f"\n🔍 STEP 2: Running Real-time Optimization Analysis...")
        print("-" * 100)
        
        try:
            optimizer = AzureCostOptimizer()
            
            # Analyze compute
            compute_recs = optimizer.optimize_compute_resources()
            
            # Analyze storage
            storage_recs = optimizer.optimize_storage_resources()
            
            # Analyze network
            network_recs = optimizer.optimize_network_resources()
            
            # Calculate totals
            total_compute_savings = sum(r['estimated_savings'] for r in compute_recs)
            total_storage_savings = sum(r['estimated_savings'] for r in storage_recs)
            total_network_savings = sum(r['estimated_savings'] for r in network_recs)
            total_savings = total_compute_savings + total_storage_savings + total_network_savings
            
            self.optimization_results = {
                'compute': compute_recs,
                'storage': storage_recs,
                'network': network_recs,
                'summary': {
                    'total_monthly_savings': total_savings,
                    'total_annual_savings': total_savings * 12,
                    'compute_savings': total_compute_savings,
                    'storage_savings': total_storage_savings,
                    'network_savings': total_network_savings
                }
            }
            
            print(f"✅ Optimization Analysis Complete")
            print(f"   Total Potential Savings: ${total_savings:.2f}/month")
            print(f"   Compute Recommendations: {len(compute_recs)}")
            print(f"   Storage Recommendations: {len(storage_recs)}")
            print(f"   Network Recommendations: {len(network_recs)}")
            
        except Exception as e:
            print(f"❌ Error in optimization analysis: {e}")
            self.optimization_results = None
    
    def create_powerpoint_report(self):
        """Create comprehensive PowerPoint report"""
        print(f"\n📝 STEP 3: Generating PowerPoint Report...")
        print("-" * 100)
        
        try:
            # Create presentation
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Slide 1: Title
            self._add_title_slide(prs)
            
            # Slide 2: Executive Summary
            self._add_executive_summary_slide(prs)
            
            # Slide 3: Cost Analysis Overview
            if self.cost_analysis_results:
                self._add_cost_analysis_slide(prs)
            
            # Slide 4: Top Services by Cost
            if self.cost_analysis_results and self.cost_analysis_results['analysis']['by_service']:
                self._add_top_services_slide(prs)
            
            # Slide 5: Optimization Summary
            if self.optimization_results:
                self._add_optimization_summary_slide(prs)
            
            # Slide 6: Compute Optimization
            if self.optimization_results and self.optimization_results['compute']:
                self._add_compute_optimization_slide(prs)
            
            # Slide 7: Storage & Network Optimization
            if self.optimization_results:
                self._add_storage_network_optimization_slide(prs)
            
            # Slide 8: Implementation Roadmap
            self._add_implementation_roadmap_slide(prs)
            
            # Save presentation
            os.makedirs('output', exist_ok=True)
            ppt_file = f'output/cost_management_report_{self.timestamp}.pptx'
            prs.save(ppt_file)
            
            print(f"✅ PowerPoint report created: {ppt_file}")
            return ppt_file
            
        except Exception as e:
            print(f"❌ Error creating PowerPoint: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def _add_title_slide(self, prs):
        """Add title slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue
        
        # Title
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = "Azure Cost Management Report"
        
        p = title_frame.paragraphs[0]
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(left, top + Inches(1.2), width, Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = f"Analysis & Optimization Report - {datetime.now().strftime('%B %d, %Y')}"
        
        p = subtitle_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    def _add_executive_summary_slide(self, prs):
        """Add executive summary slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
        
        title = slide.shapes.title
        title.text = "Executive Summary"
        
        # Get data
        if self.cost_analysis_results and self.optimization_results:
            total_cost = self.cost_analysis_results['analysis']['total_cost']
            monthly_savings = self.optimization_results['summary']['total_monthly_savings']
            annual_savings = self.optimization_results['summary']['total_annual_savings']
            
            # Add content
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4)
            
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.word_wrap = True
            
            # Current spend
            p = tf.paragraphs[0]
            p.text = f"💰 Current Monthly Spend: ${total_cost:.2f}"
            p.font.size = Pt(24)
            p.font.bold = True
            p.space_after = Pt(20)
            
            # Savings potential
            p = tf.add_paragraph()
            p.text = f"💡 Potential Monthly Savings: ${monthly_savings:.2f}"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 176, 80)  # Green
            p.space_after = Pt(20)
            
            # Annual savings
            p = tf.add_paragraph()
            p.text = f"📈 Potential Annual Savings: ${annual_savings:.2f}"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 176, 80)
            p.space_after = Pt(20)
            
            # ROI
            if total_cost > 0:
                roi = (monthly_savings / total_cost) * 100
                p = tf.add_paragraph()
                p.text = f"📊 Cost Reduction Potential: {roi:.1f}%"
                p.font.size = Pt(22)
                p.font.bold = True
    
    def _add_cost_analysis_slide(self, prs):
        """Add cost analysis slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = "Cost Analysis Overview"
        
        analysis = self.cost_analysis_results['analysis']
        days = self.cost_analysis_results.get('days_analyzed', 30)
        
        # Summary text
        left = Inches(1)
        top = Inches(2)
        width = Inches(4)
        height = Inches(4)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"Analysis Period: {days} days"
        p.font.size = Pt(18)
        p.space_after = Pt(15)
        
        p = tf.add_paragraph()
        p.text = f"Total Cost: ${analysis['total_cost']:.2f}"
        p.font.size = Pt(20)
        p.font.bold = True
        p.space_after = Pt(15)
        
        p = tf.add_paragraph()
        p.text = f"Daily Average: ${analysis['daily_average']:.2f}"
        p.font.size = Pt(18)
        p.space_after = Pt(15)
        
        p = tf.add_paragraph()
        p.text = f"Services Analyzed: {len(analysis['by_service'])}"
        p.font.size = Pt(18)
        p.space_after = Pt(15)
        
        p = tf.add_paragraph()
        p.text = f"Resource Groups: {len(analysis['by_resource_group'])}"
        p.font.size = Pt(18)
    
    def _add_top_services_slide(self, prs):
        """Add top services chart slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = "Top Services by Cost"
        
        # Get top 5 services
        services = self.cost_analysis_results['analysis']['by_service'][:5]
        
        # Create chart data
        chart_data = CategoryChartData()
        chart_data.categories = [s['name'][:20] for s in services]  # Truncate long names
        chart_data.add_series('Cost ($)', [s['cost'] for s in services])
        
        # Add chart
        x, y, cx, cy = Inches(1.5), Inches(2), Inches(7), Inches(4.5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        
        chart.has_legend = False
        chart.has_title = False
    
    def _add_optimization_summary_slide(self, prs):
        """Add optimization summary slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = "Optimization Opportunities"
        
        summary = self.optimization_results['summary']
        
        # Create pie chart data for savings breakdown
        chart_data = CategoryChartData()
        chart_data.categories = ['Compute', 'Storage', 'Network']
        chart_data.add_series('Savings', [
            summary['compute_savings'],
            summary['storage_savings'],
            summary['network_savings']
        ])
        
        # Add pie chart
        x, y, cx, cy = Inches(5), Inches(2), Inches(4.5), Inches(4.5)
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        ).chart
        
        chart.has_legend = True
        chart.has_title = False
        
        # Add summary text
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(4)
        height = Inches(4)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "Monthly Savings Breakdown"
        p.font.size = Pt(20)
        p.font.bold = True
        p.space_after = Pt(15)
        
        p = tf.add_paragraph()
        p.text = f"🖥️  Compute: ${summary['compute_savings']:.2f}"
        p.font.size = Pt(16)
        p.space_after = Pt(10)
        
        p = tf.add_paragraph()
        p.text = f"💾 Storage: ${summary['storage_savings']:.2f}"
        p.font.size = Pt(16)
        p.space_after = Pt(10)
        
        p = tf.add_paragraph()
        p.text = f"🌐 Network: ${summary['network_savings']:.2f}"
        p.font.size = Pt(16)
        p.space_after = Pt(15)
        
        p = tf.add_paragraph()
        p.text = f"Total: ${summary['total_monthly_savings']:.2f}/mo"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 176, 80)
    
    def _add_compute_optimization_slide(self, prs):
        """Add compute optimization details slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = "Compute Optimization Recommendations"
        
        compute_recs = self.optimization_results['compute']
        
        # Count by priority
        high_priority = [r for r in compute_recs if r['priority'] == 'High']
        medium_priority = [r for r in compute_recs if r['priority'] == 'Medium']
        
        left = Inches(0.5)
        top = Inches(1.8)
        width = Inches(9)
        height = Inches(5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        # Summary
        p = tf.paragraphs[0]
        p.text = f"Total VMs Analyzed: {len(compute_recs)}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.space_after = Pt(15)
        
        # High priority
        if high_priority:
            p = tf.add_paragraph()
            p.text = f"🔴 High Priority Issues: {len(high_priority)}"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(192, 0, 0)
            p.space_after = Pt(10)
            
            for rec in high_priority[:3]:  # Show top 3
                p = tf.add_paragraph()
                p.text = f"  • {rec['resource_name']}: Avg CPU {rec['avg_cpu']:.1f}% - ${rec['estimated_savings']:.0f}/mo savings"
                p.font.size = Pt(14)
                p.level = 1
                p.space_after = Pt(5)
        
        # Key recommendations
        p = tf.add_paragraph()
        p.text = "\nKey Actions:"
        p.font.size = Pt(16)
        p.font.bold = True
        p.space_after = Pt(10)
        
        actions = [
            "Purchase Reserved Instances for production VMs (up to 72% savings)",
            "Downsize severely underutilized VMs",
            "Implement auto-shutdown for dev/test environments",
            "Use B-series burstable VMs for variable workloads"
        ]
        
        for action in actions:
            p = tf.add_paragraph()
            p.text = f"  ✓ {action}"
            p.font.size = Pt(13)
            p.level = 1
            p.space_after = Pt(5)
    
    def _add_storage_network_optimization_slide(self, prs):
        """Add storage and network optimization slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = "Storage & Network Optimization"
        
        storage_recs = self.optimization_results['storage']
        network_recs = self.optimization_results['network']
        
        # Split into two columns
        # Left column - Storage
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(4.5)
        height = Inches(4.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "💾 Storage Optimization"
        p.font.size = Pt(20)
        p.font.bold = True
        p.space_after = Pt(15)
        
        p = tf.add_paragraph()
        p.text = f"Storage Accounts: {len(storage_recs)}"
        p.font.size = Pt(14)
        p.space_after = Pt(10)
        
        p = tf.add_paragraph()
        p.text = "Recommendations:"
        p.font.size = Pt(14)
        p.font.bold = True
        p.space_after = Pt(8)
        
        storage_actions = [
            "Implement lifecycle policies",
            "Move to Cool/Archive tiers",
            "Delete unattached disks",
            "Optimize redundancy settings"
        ]
        
        for action in storage_actions:
            p = tf.add_paragraph()
            p.text = f"  • {action}"
            p.font.size = Pt(12)
            p.level = 1
            p.space_after = Pt(5)
        
        # Right column - Network
        left = Inches(5.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = "🌐 Network Optimization"
        p.font.size = Pt(20)
        p.font.bold = True
        p.space_after = Pt(15)
        
        p = tf.add_paragraph()
        p.text = f"Resources Analyzed: {len(network_recs)}"
        p.font.size = Pt(14)
        p.space_after = Pt(10)
        
        p = tf.add_paragraph()
        p.text = "Recommendations:"
        p.font.size = Pt(14)
        p.font.bold = True
        p.space_after = Pt(8)
        
        network_actions = [
            "Delete unused public IPs",
            "Optimize VPN gateway SKUs",
            "Use CDN for static content",
            "Minimize cross-region traffic"
        ]
        
        for action in network_actions:
            p = tf.add_paragraph()
            p.text = f"  • {action}"
            p.font.size = Pt(12)
            p.level = 1
            p.space_after = Pt(5)
    
    def _add_implementation_roadmap_slide(self, prs):
        """Add implementation roadmap slide"""
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        title = slide.shapes.title
        title.text = "Implementation Roadmap"
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        roadmap = [
            ("Week 1-2: Quick Wins", [
                "Delete unused resources (public IPs, disks)",
                "Implement auto-shutdown for dev/test VMs",
                "Set up cost alerts and budgets"
            ]),
            ("Week 3-4: Optimization", [
                "Downsize underutilized VMs",
                "Implement storage lifecycle policies",
                "Optimize network configurations"
            ]),
            ("Month 2-3: Strategic Changes", [
                "Purchase Reserved Instances for production",
                "Implement Azure Hybrid Benefit",
                "Review and optimize database tiers"
            ]),
            ("Ongoing: Monitoring", [
                "Weekly cost reviews",
                "Monthly optimization analysis",
                "Quarterly Reserved Instance planning"
            ])
        ]
        
        for phase, actions in roadmap:
            p = tf.paragraphs[0] if phase == roadmap[0][0] else tf.add_paragraph()
            p.text = phase
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 51, 102)
            p.space_after = Pt(8)
            
            for action in actions:
                p = tf.add_paragraph()
                p.text = f"  ✓ {action}"
                p.font.size = Pt(13)
                p.level = 1
                p.space_after = Pt(4)
            
            if phase != roadmap[-1][0]:
                tf.add_paragraph().space_after = Pt(10)
    
    def save_json_summary(self):
        """Save combined JSON summary"""
        os.makedirs('output', exist_ok=True)
        
        summary_file = f'output/cost_management_summary_{self.timestamp}.json'
        
        summary = {
            'timestamp': datetime.now().isoformat(),
            'cost_analysis': self.cost_analysis_results,
            'optimization': self.optimization_results
        }
        
        with open(summary_file, 'w') as f:
            json.dump(summary, f, indent=2, default=str)
        
        print(f"✅ JSON summary saved: {summary_file}")
        return summary_file
    
    def run_complete_analysis(self, days: int = 30):
        """Run complete analysis and generate reports"""
        print(f"\n{'='*100}")
        print("AZURE COST MANAGEMENT - COMPLETE ANALYSIS & REPORTING")
        print(f"{'='*100}\n")
        
        # Step 1: Cost Analysis
        self.run_cost_analysis(days)
        
        # Step 2: Optimization Analysis
        self.run_optimization_analysis()
        
        # Step 3: Generate PowerPoint
        ppt_file = self.create_powerpoint_report()
        
        # Step 4: Save JSON summary
        json_file = self.save_json_summary()
        
        # Final summary
        print(f"\n{'='*100}")
        print("✅ COMPLETE ANALYSIS FINISHED")
        print(f"{'='*100}")
        
        if self.cost_analysis_results and self.optimization_results:
            print(f"\n📊 Results Summary:")
            print(f"   Current Monthly Spend: ${self.cost_analysis_results['analysis']['total_cost']:.2f}")
            print(f"   Potential Monthly Savings: ${self.optimization_results['summary']['total_monthly_savings']:.2f}")
            print(f"   Potential Annual Savings: ${self.optimization_results['summary']['total_annual_savings']:.2f}")
            
            if self.cost_analysis_results['analysis']['total_cost'] > 0:
                roi = (self.optimization_results['summary']['total_monthly_savings'] / 
                       self.cost_analysis_results['analysis']['total_cost']) * 100
                print(f"   Cost Reduction Potential: {roi:.1f}%")
        
        print(f"\n📄 Generated Files:")
        if ppt_file:
            print(f"   ✅ PowerPoint: {ppt_file}")
        if json_file:
            print(f"   ✅ JSON Summary: {json_file}")
        
        print(f"\n{'='*100}\n")


def main():
    """Main entry point"""
    orchestrator = CostManagementOrchestrator()
    orchestrator.run_complete_analysis(days=30)


if __name__ == "__main__":
    main()
